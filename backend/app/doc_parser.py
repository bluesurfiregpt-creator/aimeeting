"""
Document text extractor.

Inputs: filename + raw bytes. Output: a single plaintext string.

Format coverage for Sprint I:
- pdf  : pypdf (fast, lightweight, no system deps)
- docx : python-docx
- xlsx : openpyxl, joining each row's cells with tab + newline between rows
- txt / md / text-shaped extensions : decode as UTF-8 (with fallback)

Anything else → ValueError so the caller can mark the document failed.
"""

from __future__ import annotations

import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)


SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    # v27.0-mobile P19-B.2: PPTX 抽 — slide 标题 + 内容 文本 框 (无图)
    ".pptx": "pptx",
    ".txt": "text",
    ".md": "text",
    ".markdown": "text",
    ".text": "text",
    ".csv": "text",
    ".log": "text",
    ".json": "text",
    ".yaml": "text",
    ".yml": "text",
    # v25-2: 图片格式 → OCR(Qwen-VL).政务客户大量历史扫描件 / 拍照公文
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".bmp": "image",
    ".tiff": "image",
    ".tif": "image",
    ".webp": "image",
    ".gif": "image",
}


def kind_from_filename(filename: str) -> Optional[str]:
    name = (filename or "").lower()
    for ext, kind in SUPPORTED_EXTENSIONS.items():
        if name.endswith(ext):
            return kind
    return None


def extract_text(filename: str, content: bytes) -> str:
    """
    同步抽取 — 用于不需要 OCR 的格式.

    image / 扫描件 PDF 走 async 路径(extract_text_async).
    若 caller 想统一 sync 接口,扫描 PDF 抽出空字符串后由 caller 决定要不要
    重试为 async 路径.
    """
    kind = kind_from_filename(filename)
    if kind is None:
        raise ValueError(f"unsupported file type: {filename}")
    if kind == "image":
        # 同步路径不能调 OCR(async).让 caller 走 extract_text_async.
        raise ValueError(
            f"image file '{filename}' requires async OCR path "
            f"(use extract_text_async instead of extract_text)"
        )
    if kind == "pdf":
        return _extract_pdf(content)
    if kind == "docx":
        return _extract_docx(content)
    if kind == "xlsx":
        return _extract_xlsx(content)
    if kind == "pptx":
        return _extract_pptx(content)
    if kind == "text":
        return _decode_text(content)
    raise ValueError(f"unhandled kind: {kind}")


async def extract_text_async(filename: str, content: bytes) -> str:
    """
    v25-2 异步抽取(支持 OCR).

    行为:
      - 图片 → 直接 OCR
      - PDF → 先 pypdf;若文字量 < 阈值,自动 OCR fallback(扫描件)
      - 其他格式 → 同 extract_text(sync)

    OCR 失败时:对图片 raise OCRError(没有兜底);对 PDF 返回 pypdf 已抽到
    的文字(可能空)— 不破坏 caller 期望.
    """
    kind = kind_from_filename(filename)
    if kind is None:
        raise ValueError(f"unsupported file type: {filename}")
    if kind == "image":
        from .ocr import image_mime_for_filename, ocr_image
        return await ocr_image(content, mime_type=image_mime_for_filename(filename))
    if kind == "pdf":
        # 先 pypdf 抽 — 文字 PDF 这一步就够了,省 OCR 钱
        text = _extract_pdf(content)
        if text and len(text.strip()) >= 30:
            return text
        # 文字太少 → 大概率是扫描件 → OCR fallback
        from .ocr import maybe_ocr_pdf_if_empty
        final_text, _ = await maybe_ocr_pdf_if_empty(text, content)
        return final_text
    # 非 PDF / 非图片 → 复用 sync 路径
    return extract_text(filename, content)


def _decode_text(content: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    # Last resort
    return content.decode("utf-8", errors="replace")


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception:
            logger.exception("pdf page %d extraction failed", i)
            text = ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def _extract_docx(content: bytes) -> str:
    """
    DOCX 文本提取.两段式:
      1. 主路径 python-docx — 表格保留完好,段落整齐
      2. fallback docx2txt — 主路径炸时(WPS / 在线 Office 导出经常缺
         word/endnotes.xml / footnotes.xml 等可选 part,python-docx
         严格 KeyError 拒绝),用 docx2txt 强行抽文本.精度低一点
         (表格变线性,部分格式丢)但**不会失败**.

    实测最常见 fallback 触发原因:WPS / iWork / Google Docs 导出的 .docx
    缺 endnotes.xml,以及 OnlyOffice 缺 footer 引用.
    """
    try:
        return _extract_docx_strict(content)
    except KeyError as e:
        # zipfile.ZipFile.open() KeyError 长成 "There is no item named '...' in the archive"
        logger.warning("python-docx 失败(%s),fallback 到 docx2txt", e)
        return _extract_docx_loose(content)
    except Exception as e:
        # 其他异常(BadZipFile / 损坏 / 加密 等)也走 fallback,docx2txt 也会
        # 抛 — 那就让它抛上去,让 router 标 status=failed.
        logger.warning("python-docx 异常(%s: %s),fallback 到 docx2txt", type(e).__name__, e)
        return _extract_docx_loose(content)


def _extract_docx_strict(content: bytes) -> str:
    """python-docx 主路径 — 表格 + 段落都漂亮."""
    from docx import Document  # python-docx

    doc = Document(io.BytesIO(content))
    parts: list[str] = []
    # Body paragraphs
    for p in doc.paragraphs:
        if p.text and p.text.strip():
            parts.append(p.text)
    # Tables (joined cell text per row, tab-separated)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_docx_loose(content: bytes) -> str:
    """
    docx2txt fallback — 把 docx 当 zip 打开,抽 word/document.xml 的文本.
    需要写到临时文件(docx2txt API 只接 path).
    """
    import tempfile
    import os

    import docx2txt

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name
    try:
        text = docx2txt.process(tmp_path) or ""
        return text.strip()
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"### {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if v is None else str(v).strip() for v in row]
            # Skip rows that are all empty
            if any(c for c in cells):
                parts.append("\t".join(cells))
        parts.append("")
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    """v27.0-mobile P19-B.2: 抽 .pptx 文字内容 (slide 标题 + 文本框 + notes).

    每 slide 用 "### slide N" header 分隔; 同 slide 内 标题 / 正文 各占 一行.
    备注 (presenter notes) 也 抽 — 政务 PPT 经常 把 解读 写 notes 里.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError("python-pptx 未安装 — 添加到 requirements.txt 后重启")

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"### slide {i}")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = getattr(shape, "text_frame", None)
            if tf is None:
                continue
            for para in tf.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    parts.append(line)
        # presenter notes
        if slide.has_notes_slide and slide.notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf and notes_tf.text and notes_tf.text.strip():
                parts.append(f"[备注] {notes_tf.text.strip()}")
        parts.append("")
    return "\n".join(parts).strip()
